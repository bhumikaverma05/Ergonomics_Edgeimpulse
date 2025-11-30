"""
posture_app.py
Complete AI Posture Correction desktop application:
- Uses webcam + Mediapipe to detect posture
- Tkinter UI with video preview, status, controls
- Live matplotlib graph of good/bad posture time
- CSV logging for session
- Stretching suggestions
- Export session summary to PPTX (optional: requires python-pptx)
"""

import threading
import time
import csv
import os
from datetime import datetime
import tkinter as tk
from tkinter import ttk, messagebox, filedialog

# OpenCV / Mediapipe / numpy
import cv2
import mediapipe as mp
import numpy as np

# plotting in Tkinter
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg

# For images in Tkinter
from PIL import Image, ImageTk

# Optional: winsound for Windows beep
try:
    import winsound
    def beep(freq=1000, duration=600):
        winsound.Beep(freq, duration)
except Exception:
    # fallback beep (cross-platform minimal)
    import sys
    def beep(freq=1000, duration=600):
        print("\a", end="", flush=True)

# Optional python-pptx: used only when user clicks 'Export PPT'
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# ---------------------------
# Posture detection utilities
# ---------------------------
mp_pose = mp.solutions.pose
mp_draw = mp.solutions.drawing_utils

def calculate_angle(a, b, c):
    a = np.array(a)
    b = np.array(b)
    c = np.array(c)
    radians = np.arctan2(c[1]-b[1], c[0]-b[0]) - np.arctan2(a[1]-b[1], a[0]-b[0])
    angle = np.abs(radians * 180.0 / np.pi)
    if angle > 180.0:
        angle = 360 - angle
    return angle

# ---------------------------
# Main App Class
# ---------------------------
class PostureApp:
    def __init__(self, root):
        self.root = root
        root.title("AI Posture Correction System — Bhumika")
        self.running = False

        # Video capture & mediapipe
        self.cap = None
        self.pose = mp_pose.Pose(min_detection_confidence=0.5, min_tracking_confidence=0.5)

        # Timers & stats
        self.bad_posture_start = 0
        self.alert_given = False
        self.good_time = 0.0
        self.bad_time = 0.0
        self.last_update_time = time.time()
        self.session_start = time.time()

        # Graph data
        self.times = []
        self.good_series = []
        self.bad_series = []

        # Posture threshold (angle); adjustable through UI
        self.angle_threshold = tk.DoubleVar(value=150.0)
        self.alert_delay = tk.DoubleVar(value=5.0)  # seconds

        # Build UI
        self.build_ui()

        # Stretching suggestions
        self.stretches = [
            ("Neck tilt", "Gently tilt your head to one side and hold for 15-20s, then switch."),
            ("Chin tuck", "Pull chin slightly backwards to align the neck, hold 10s, repeat 5 times."),
            ("Shoulder rolls", "Roll shoulders back 10 times to open chest and release tension."),
            ("Seated twist", "Rotate upper body while seated, hold 10s each side."),
            ("Chest stretch", "Clasp hands behind back and lift gently to open chest.")
        ]

    def build_ui(self):
        # Left frame: video + controls
        left = ttk.Frame(self.root)
        left.grid(row=0, column=0, padx=8, pady=8, sticky="ns")

        # Video panel (label will hold frames)
        self.video_label = ttk.Label(left)
        self.video_label.grid(row=0, column=0, columnspan=3)

        # Controls
        ttk.Button(left, text="Start", command=self.start).grid(row=1, column=0, pady=6, sticky="ew")
        ttk.Button(left, text="Stop", command=self.stop).grid(row=1, column=1, pady=6, sticky="ew")
        ttk.Button(left, text="Export PPTX", command=self.export_pptx).grid(row=1, column=2, pady=6, sticky="ew")

        # Threshold slider + delay
        ttk.Label(left, text="Angle Threshold (lower = stricter)").grid(row=2, column=0, columnspan=3, sticky="w")
        ttk.Scale(left, from_=120, to=175, variable=self.angle_threshold, orient="horizontal").grid(row=3, column=0, columnspan=3, sticky="ew")
        ttk.Label(left, text="Alert Delay (sec)").grid(row=4, column=0, columnspan=3, sticky="w")
        ttk.Scale(left, from_=1, to=10, variable=self.alert_delay, orient="horizontal").grid(row=5, column=0, columnspan=3, sticky="ew")

        # Status labels
        self.status_var = tk.StringVar(value="Status: Idle")
        self.angle_var = tk.StringVar(value="Angle: -")
        self.timer_var = tk.StringVar(value="Bad Timer: 0.0s")
        ttk.Label(left, textvariable=self.status_var).grid(row=6, column=0, columnspan=3, sticky="w")
        ttk.Label(left, textvariable=self.angle_var).grid(row=7, column=0, columnspan=3, sticky="w")
        ttk.Label(left, textvariable=self.timer_var).grid(row=8, column=0, columnspan=3, sticky="w")

        # Stretch suggestions box
        ttk.Label(left, text="Stretch Suggestions:").grid(row=9, column=0, columnspan=3, sticky="w")
        self.suggestion_box = tk.Text(left, width=40, height=6, wrap="word")
        self.suggestion_box.grid(row=10, column=0, columnspan=3, pady=4)
        self.suggestion_box.insert("end", "Good posture tips will appear here when bad posture persists.")
        self.suggestion_box.config(state="disabled")

        # Right frame: Graph
        right = ttk.Frame(self.root)
        right.grid(row=0, column=1, padx=8, pady=8)

        self.fig = Figure(figsize=(5,3), dpi=100)
        self.ax = self.fig.add_subplot(111)
        self.ax.set_title("Good vs Bad Time (seconds)")
        self.ax.set_xlabel("Time (s)")
        self.ax.set_ylabel("Seconds")
        self.line_good, = self.ax.plot([], [], label="Good")
        self.line_bad, = self.ax.plot([], [], label="Bad")
        self.ax.legend()

        self.canvas = FigureCanvasTkAgg(self.fig, master=right)
        self.canvas.get_tk_widget().pack()

    # ---------------------------
    # Start / Stop
    # ---------------------------
    def start(self):
        if self.running:
            return
        self.running = True
        self.status_var.set("Status: Running")
        self.cap = cv2.VideoCapture(0, cv2.CAP_DSHOW)  # CAP_DSHOW helps on Windows
        self.session_start = time.time()
        self.last_update_time = time.time()
        self.loop_thread = threading.Thread(target=self.video_loop, daemon=True)
        self.loop_thread.start()
        self.update_graph_loop()

    def stop(self):
        if not self.running:
            return
        self.running = False
        self.status_var.set("Status: Stopping...")
        # Wait briefly for thread to stop
        time.sleep(0.3)
        if self.cap:
            self.cap.release()
        self.status_var.set("Status: Idle")
        # Save CSV summary
        self.save_csv()
        messagebox.showinfo("Session saved", "Session saved to posture_log.csv")

    # ---------------------------
    # Video / Detection loop
    # ---------------------------
    def video_loop(self):
        with mp_pose.Pose(min_detection_confidence=0.5, min_tracking_confidence=0.5) as pose:
            while self.running:
                success, frame = self.cap.read()
                if not success:
                    continue

                # Flip for mirror effect
                frame = cv2.flip(frame, 1)
                h, w, _ = frame.shape

                # Convert color
                img_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                results = pose.process(img_rgb)

                angle = None
                if results.pose_landmarks:
                    lm = results.pose_landmarks.landmark
                    # Use LEFT shoulder/ear/hip (works if facing camera)
                    try:
                        shoulder = [lm[mp_pose.PoseLandmark.LEFT_SHOULDER.value].x * w,
                                    lm[mp_pose.PoseLandmark.LEFT_SHOULDER.value].y * h]
                        ear = [lm[mp_pose.PoseLandmark.LEFT_EAR.value].x * w,
                               lm[mp_pose.PoseLandmark.LEFT_EAR.value].y * h]
                        hip = [lm[mp_pose.PoseLandmark.LEFT_HIP.value].x * w,
                               lm[mp_pose.PoseLandmark.LEFT_HIP.value].y * h]
                        angle = calculate_angle(ear, shoulder, hip)
                        mp_draw.draw_landmarks(frame, results.pose_landmarks, mp_pose.POSE_CONNECTIONS)
                    except Exception:
                        angle = None

                # Timing & stats
                current_time = time.time()
                time_diff = current_time - self.last_update_time
                self.last_update_time = current_time

                if angle is not None:
                    self.angle_var.set(f"Angle: {int(angle)}")
                    threshold = self.angle_threshold.get()
                    delay_needed = self.alert_delay.get()

                    if angle < threshold:
                        # BAD posture
                        self.bad_time += time_diff
                        # start timer if not started
                        if self.bad_posture_start == 0:
                            self.bad_posture_start = time.time()
                        elapsed_bad = time.time() - self.bad_posture_start
                        self.timer_var.set(f"Bad Timer: {elapsed_bad:.1f}s")

                        # Provide beep & suggestion after delay
                        if elapsed_bad > delay_needed and not self.alert_given:
                            self.alert_given = True
                            # beep (non-blocking)
                            threading.Thread(target=beep, args=(1000, 600), daemon=True).start()
                            # show suggestion
                            self.show_suggestion()
                    else:
                        # GOOD posture
                        self.good_time += time_diff
                        self.bad_posture_start = 0
                        self.alert_given = False
                        self.timer_var.set("Bad Timer: 0.0s")
                else:
                    self.angle_var.set("Angle: -")

                # Update graph data arrays
                elapsed_session = current_time - self.session_start
                self.times.append(elapsed_session)
                self.good_series.append(self.good_time)
                self.bad_series.append(self.bad_time)

                # Convert to PhotoImage and update Tkinter label
                img_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                img_pil = Image.fromarray(img_rgb).resize((640, 480))
                imgtk = ImageTk.PhotoImage(image=img_pil)
                # keep a reference to avoid GC
                self.video_label.imgtk = imgtk
                self.video_label.config(image=imgtk)

            # release handled by stop()

    # ---------------------------
    # Graph update loop (runs on main thread)
    # ---------------------------
    def update_graph_loop(self):
        if not self.running:
            return
        # update plot
        self.line_good.set_data(self.times, self.good_series)
        self.line_bad.set_data(self.times, self.bad_series)
        # adjust axes
        if self.times:
            self.ax.set_xlim(max(0, self.times[-1]-60), max(10, self.times[-1]+1))
            ymax = max(max(self.good_series, default=1), max(self.bad_series, default=1), 1)
            self.ax.set_ylim(0, ymax*1.2)
        self.canvas.draw_idle()
        # schedule next update
        self.root.after(1000, self.update_graph_loop)

    # ---------------------------
    # Suggestions UI
    # ---------------------------
    def show_suggestion(self):
        # choose a suggestion (rotate based on bad_time)
        idx = int(self.bad_time) % len(self.stretches)
        title, desc = self.stretches[idx]
        text = f"{title}\n\n{desc}\n\n(Do this 2-3 times and correct posture.)"
        self.suggestion_box.config(state="normal")
        self.suggestion_box.delete("1.0", "end")
        self.suggestion_box.insert("end", text)
        self.suggestion_box.config(state="disabled")

    # ---------------------------
    # Save CSV
    # ---------------------------
    def save_csv(self):
        filename = "posture_log.csv"
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(filename, "a", newline="") as f:
            writer = csv.writer(f)
            writer.writerow([now,
                             f"Good_Time_s={round(self.good_time,2)}",
                             f"Bad_Time_s={round(self.bad_time,2)}",
                             f"Threshold={self.angle_threshold.get()}",
                             f"Delay_s={self.alert_delay.get()}"])
        print("Saved session to", filename)

    # ---------------------------
    # Export PPTX summary (optional)
    # ---------------------------
    def export_pptx(self):
        if not PPTX_AVAILABLE:
            messagebox.showinfo("python-pptx missing",
                                "python-pptx is not installed. Install it with:\n\npip install python-pptx")
            return

        # Create presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Posture Session Summary"
        subtitle.text = f"Recorded: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # Add stats slide
        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Session Stats"
        body = slide2.shapes.placeholders[1].text_frame
        body.text = f"Good time: {round(self.good_time,2)} s"
        p = body.add_paragraph()
        p.text = f"Bad time: {round(self.bad_time,2)} s"
        p.level = 1
        p = body.add_paragraph()
        p.text = f"Angle threshold: {self.angle_threshold.get()}"
        p.level = 1
        p = body.add_paragraph()
        p.text = f"Alert delay (s): {self.alert_delay.get()}"
        p.level = 1

        # Save file
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                 filetypes=[("PowerPoint Files", "*.pptx")],
                                                 title="Save presentation as")
        if save_path:
            prs.save(save_path)
            messagebox.showinfo("Saved", f"PPTX saved to {save_path}")

# ---------------------------
# Run the app
# ---------------------------
def main():
    root = tk.Tk()
    app = PostureApp(root)

    # handle closing
    def on_close():
        if app.running:
            if messagebox.askyesno("Quit", "App is running. Stop and exit?"):
                app.stop()
            else:
                return
        root.destroy()
    root.protocol("WM_DELETE_WINDOW", on_close)
    root.mainloop()

if __name__ == "__main__":
    main()
